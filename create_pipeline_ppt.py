#!/usr/bin/env python3
"""Generate the Korean T9-GPT implementation overview deck."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "T9-GPT-Pipeline-Implementation-KR-Fixed.pptx"

BG = RGBColor(0x0D, 0x11, 0x17)
PANEL = RGBColor(0x16, 0x1B, 0x22)
PANEL_2 = RGBColor(0x21, 0x26, 0x2D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0xC9, 0xD1, 0xD9)
MUTED = RGBColor(0x8B, 0x94, 0x9E)
RED = RGBColor(0xE0, 0x30, 0x30)
CYAN = RGBColor(0x00, 0xD4, 0xFF)
GREEN = RGBColor(0x3F, 0xB9, 0x50)
YELLOW = RGBColor(0xF7, 0xC9, 0x48)
PURPLE = RGBColor(0xA3, 0x71, 0xF7)
BORDER = RGBColor(0x30, 0x36, 0x3D)

FONT = "Malgun Gothic"
MONO = "Courier New"
SW = Inches(13.333)
SH = Inches(7.5)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 14,
    color: RGBColor = TEXT,
    bold: bool = False,
    font: str = FONT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    margin: float = 0.03,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    for idx, line in enumerate(text.split("\n")):
        if idx:
            p = tf.add_paragraph()
            p.alignment = align
            p.space_after = Pt(0)
            p.space_before = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return box


def add_rich_line(slide, runs, x, y, w, h, *, size=14, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    for text, color, bold, font in runs:
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return box


def rect(slide, x, y, w, h, fill=PANEL, line=BORDER, radius=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def pill(slide, text, x, y, w, color, *, size=10.5):
    shape = rect(slide, x, y, w, 0.32, fill=color, line=color)
    add_text(
        slide,
        text,
        x,
        y + 0.01,
        w,
        0.28,
        size=size,
        color=BG,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def title(slide, heading, subtitle, index, total):
    add_text(slide, heading, 0.35, 0.08, 12.6, 0.55, size=29, color=WHITE, bold=True)
    add_text(slide, subtitle, 0.35, 0.72, 12.6, 0.35, size=13.5, color=CYAN)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(1.17), Inches(12.62), Inches(0.025)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER
    line.line.fill.background()
    add_text(
        slide,
        f"t9project.dev  ·  T9-GPT  ·  {index}/{total}",
        0.35,
        7.13,
        12.62,
        0.25,
        size=9.5,
        color=MUTED,
    )


def bullet_list(slide, items, x, y, w, h, *, size=13.5, color=TEXT, bullet_color=CYAN):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(9)
        p.level = 0
        r1 = p.add_run()
        r1.text = "▸  "
        r1.font.name = FONT
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = bullet_color
        r2 = p.add_run()
        r2.text = item
        r2.font.name = FONT
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return box


def section_label(slide, text, x, y, w, color=CYAN):
    add_text(slide, text, x, y, w, 0.34, size=15.5, color=color, bold=True)


def code_block(slide, text, x, y, w, h, *, size=10.5):
    rect(slide, x, y, w, h, fill=RGBColor(0x08, 0x0B, 0x10), line=BORDER)
    return add_text(slide, text, x + 0.15, y + 0.13, w - 0.3, h - 0.2, size=size, color=GREEN, font=MONO)


def card(slide, heading, body, x, y, w, h, color=CYAN, *, body_size=12.5):
    rect(slide, x, y, w, h)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.06), Inches(h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(slide, heading, x + 0.18, y + 0.14, w - 0.3, 0.34, size=14, color=color, bold=True)
    add_text(slide, body, x + 0.18, y + 0.58, w - 0.3, h - 0.7, size=body_size, color=TEXT)


def arrow(slide, x, y, w=0.42, h=0.32, color=MUTED):
    return add_text(
        slide,
        "→",
        x,
        y - 0.04,
        w,
        h + 0.08,
        size=18,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0,
    )


def connector(slide, x1, y1, x2, y2, color=MUTED, width=1.5):
    thickness = max(0.012, width / 72)
    x = min(x1, x2)
    y = min(y1, y2)
    w = max(abs(x2 - x1), thickness)
    h = max(abs(y2 - y1), thickness)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    blank = prs.slide_layouts[6]
    total = 13

    # 1. Cover
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    add_text(s, "T9", 0.42, 0.82, 3.3, 1.65, size=118, color=RED, bold=True)