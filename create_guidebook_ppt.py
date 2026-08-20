#!/usr/bin/env python3
"""Generate the Korean T9-GPT beginner's guidebook deck.

Written for someone who has never seen T9-GPT before: what problem it solves,
how the (now single-mode, claude-loop-only) pipeline actually works end to
end, how a scenario file is put together, how success is judged, what safety
rails exist, and what ends up in the output directory. Reuses the visual
style/helpers from create_pipeline_ppt.py so it matches the rest of the deck
family, but every slide's content is new and reflects the current, cleaned-up
codebase (no deterministic run mode, no PentestGPT engine).
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from create_pipeline_ppt import (
    BG, BORDER, CYAN, GREEN, MONO, MUTED, PANEL, PANEL_2, PURPLE, RED, TEXT,
    WHITE, YELLOW, SH, SW,
    add_text, arrow, bullet_list, card, code_block, pill, rect,
    section_label, title,
)

OUT = Path(__file__).resolve().parent.parent / "T9-GPT-Guidebook-KR.pptx"
TOTAL = 17


def _bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    blank = prs.slide_layouts[6]

    # 1. Cover ------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    add_text(s, "T9", 0.42, 0.7, 3.3, 1.65, size=110, color=RED, bold=True)
    add_text(s, "GPT", 2.78, 1.3, 4.2, 1.15, size=76, color=WHITE, bold=True)
    add_text(s, "가이드북 — 처음 보는 사람을 위한 개요", 0.45, 3.1, 12.0, 0.65, size=24, color=TEXT, bold=True)
    add_text(
        s,
        "사람이 표적(CVE)만 지정 → AI 에이전트가 스스로 침투 → 실제 트래픽으로 성공 여부를 검증 → 라벨링된 데이터셋",
        0.45, 3.82, 12.4, 0.4, size=13.5, color=CYAN,
    )
    pill(s, "GUIDEBOOK", 0.45, 4.55, 1.9, GREEN)
    pill(s, "CLAUDE-LOOP 단일 엔진", 2.5, 4.55, 2.55, CYAN)
    pill(s, "DOCKER / VULHUB", 5.2, 4.55, 2.15, YELLOW)
    pill(s, "NDR · PCAP 데이터셋", 7.5, 4.55, 2.35, PURPLE)
    add_text(
        s,
        "이 문서는 IDS/NDR/EDR 탐지 모델 학습용 공격 트래픽을 자동 생성하는 T9-GPT의\n"
        "구조·동작 원리·사용법을 처음 보는 사람도 이해할 수 있도록 정리한 가이드북입니다.",
        0.45, 5.3, 11.5, 0.85, size=12.5, color=MUTED,
    )
    add_text(s, "2026", 11.65, 6.46, 1.2, 0.5, size=18, color=RED, bold=True, align=PP_ALIGN.RIGHT)
    add_text(s, "t9project.dev", 0.45, 6.68, 4.0, 0.3, size=13, color=CYAN)

    # 2. What / Why ---------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "T9-GPT란 무엇인가", "침입 탐지 AI를 학습시키려면 '라벨이 붙은 공격 트래픽'이 많이 필요하다", 2, TOTAL)
    card(s, "① 필요한 것", "IDS·NDR·EDR 같은 침입 탐지 모델을\n학습시키려면 진짜 공격 트래픽 +\n정확한 라벨(공격/정상, 기법)이 필요",
         0.45, 1.5, 3.95, 2.0, CYAN)
    card(s, "② 기존 방식의 한계", "사람이 CVE마다 익스플로잇 코드를\n직접 작성 → 느리고, 같은 스크립트는\n항상 같은 트래픽만 생성",
         4.55, 1.5, 3.95, 2.0, RED)
    card(s, "③ T9-GPT의 해법", "표적과 CVE만 정해주면 AI 에이전트가\n스스로 다른 방식으로 침투를 시도\n→ 다양한 트래픽이 자동으로 쌓임",
         8.65, 1.5, 4.22, 2.0, GREEN)
    rect(s, 0.45, 3.85, 12.43, 2.4, fill=PANEL)
    section_label(s, "한 줄 요약", 0.7, 4.1, 4.0, YELLOW)
    bullet_list(
        s,
        [
            "대상: Docker/Vulhub로 재현 가능한, 이미 공개된(CVE가 알려진) 웹 취약점.",
            "에이전트는 Claude CLI 기반 'claude-loop' — 우리가 시스템 프롬프트·도구·예산을 완전히 통제한다.",
            "성공 여부는 에이전트의 말이 아니라 실제 패킷/서버 상태(ground truth)로 판정한다.",
            "실행이 끝나면 컨테이너·네트워크를 완전히 정리하고, 라벨이 붙은 샘플 하나를 남긴다.",
        ],
        0.7, 4.5, 11.9, 1.7, size=13.5, bullet_color=YELLOW,
    )

    # 3. Core idea ------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "핵심 아이디어 — 통제된 자유", "사람이 범위를 고정하고, 에이전트는 그 안에서만 방법을 자유롭게 정한다", 3, TOTAL)
    rect(s, 0.48, 1.5, 6.0, 4.6, fill=PANEL, line=GREEN)
    section_label(s, "사람이 고정하는 것", 0.75, 1.78, 5.5, GREEN)
    bullet_list(
        s,
        [
            "공격할 표적(CVE)과 Vulhub 환경",
            "쓸 수 있는 익스플로잇 '기법 목록'(기법 뱅크)",
            "실행당 달러 예산 상한 (--max-budget-usd)",
            "성공을 어떻게 판정할지(proof oracle)",
        ],
        0.75, 2.3, 5.5, 3.5, size=13.5, bullet_color=GREEN,
    )
    arrow(s, 6.62, 3.7, 0.6, 0.5, CYAN)
    rect(s, 7.4, 1.5, 5.48, 4.6, fill=PANEL, line=CYAN)
    section_label(s, "에이전트가 자유롭게 정하는 것", 7.67, 1.78, 5.0, CYAN)
    bullet_list(
        s,
        [
            "정해진 기법 안에서 정확한 payload·인코딩",
            "curl 요청의 구체적인 형태와 순서",
            "실행마다 다른 기법을 배정받아 다른 트래픽 생성",
            "재현용 exploit 스크립트를 어떻게 작성할지",
        ],
        7.67, 2.3, 5.0, 3.5, size=13.5, bullet_color=CYAN,
    )
    add_text(
        s,
        "→ 같은 CVE를 대상으로 하되, 실행마다 조금씩 다른 방식으로 침투한다 — '다양하지만 주제에서 벗어나지 않는' 데이터",
        0.6, 6.3, 12.2, 0.45, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
    )

    # 4. Module map ------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "전체 구조 한눈에 — 파일별 역할", "각 파이썬 파일이 하나의 역할만 맡는다 (단일 모드: 에이전트 파이프라인)", 4, TOTAL)
    mods = [
        ("orchestrator.py", "전체 실행을 지휘하는 CLI. list·validate·agent 3개 명령을 제공하고, 한 번의 실행(타겟 기동→캡처→에이전트→판정→정리)을 끝까지 관리한다.", CYAN),
        ("models.py", "시나리오 JSON의 설계도(Pydantic 스키마). 잘못된 값·위험한 경로를 사전에 걸러내고, 기법 뱅크·예산·판정 방식을 정의한다.", GREEN),
        ("environment.py", "VulhubTarget 클래스로 Docker Compose 타겟을 켜고(compose up), 준비될 때까지 기다리고, 끝나면 반드시 끈다(teardown).", YELLOW),
        ("agent_runner.py", "실제 claude-loop 에이전트를 실행하는 엔진. 시스템 프롬프트·허용 도구·예산 상한을 강제하고 결과를 파싱한다.", PURPLE),
        ("collector.py", "tcpdump로 네트워크 트래픽(PCAP)을 캡처하고, proof token이 응답/로그/파일 어디에 나타났는지 검사(oracle)한다.", RED),
        ("scenarios/*.json", "사람이 직접 작성하는 설정 파일 — 어떤 CVE를, 어떤 기법으로, 얼마의 예산으로 공격할지를 정의한다.", WHITE),
    ]
    y = 1.5
    for name, desc, color in mods:
        rect(s, 0.45, y, 12.43, 0.82, fill=PANEL, line=BORDER)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(y), Inches(0.06), Inches(0.82))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_text(s, name, 0.66, y + 0.1, 2.55, 0.62, size=13, color=color, bold=True, font=MONO, valign=1)
        add_text(s, desc, 3.35, y + 0.08, 9.35, 0.68, size=11.3, color=TEXT)
        y += 0.9

    # 5. Commands ---------------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "실행 명령어 — 딱 3가지", "list · validate · agent — 데이터를 만드는 건 agent 명령 하나뿐", 5, TOTAL)
    section_label(s, "① 카탈로그 확인", 0.47, 1.4, 6.0, CYAN)
    code_block(s, "uv run python orchestrator.py list \\\n  --config scenarios/example.json", 0.47, 1.8, 6.0, 0.95, size=11.5)
    add_text(s, "등록된 시나리오와 실행 가능 여부를 표로 보여준다.", 0.47, 2.85, 6.0, 0.4, size=11.5, color=MUTED)
    section_label(s, "② 카탈로그 검증", 6.9, 1.4, 6.0, GREEN)
    code_block(s, "uv run python orchestrator.py validate \\\n  --config scenarios/example.json", 6.9, 1.8, 6.0, 0.95, size=11.5)
    add_text(s, "형식이 잘못됐거나 실행 불가능한 시나리오가 있으면 실패로 종료한다.", 6.9, 2.85, 6.0, 0.4, size=11.5, color=MUTED)
    section_label(s, "③ 데이터 수집 — 실제 실행", 0.47, 3.55, 12.0, YELLOW)
    code_block(
        s,
        "uv run python orchestrator.py agent --config scenarios/example.json \\\n"
        "  --t9-code T9-25-02-S-N-CD [--benign] [--allow-web] [--repeat N]\\\n"
        "  [--technique ID] [--model NAME] [--budget USD] [--attempts N]",
        0.47, 3.95, 12.4, 1.0, size=11.3,
    )
    bullet_list(
        s,
        [
            "--benign : 공격 대신 '정상 사용자처럼 보이는' 트래픽 샘플을 수집한다.",
            "--repeat N : 같은 시나리오를 N번 실행 — 기법·모델이 자동으로 로테이션되어 다양성이 생긴다.",
            "--allow-web : 에이전트에게 웹 검색/조회 권한을 준다 (공개된 CVE 조사 전용 — 뒤에서 다룸).",
        ],
        0.47, 5.1, 12.2, 1.3, size=12.5, bullet_color=YELLOW,
    )

    # 6. Lifecycle ----------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "한 번의 실행이 흘러가는 과정", "토큰(비용)은 딱 한 단계, '에이전트 실행'에서만 소비된다", 6, TOTAL)
    steps = [
        ("1", "타겟 기동", "Vulhub Docker\ncompose up", CYAN),
        ("2", "참고자료 주입", "선택 사항\n(RAG)", GREEN),
        ("3", "baseline 점검", "선택 사항\n무료", YELLOW),
        ("4", "캡처 시작", "tcpdump\n트래픽 기록", RED),
        ("5", "에이전트 실행", "비용 발생\n지점", PURPLE),
        ("6", "Ground-Truth 판정", "proof oracle\n로 라벨링", CYAN),
        ("7", "완전 정리", "teardown\n항상 실행", WHITE),
    ]
    x = 0.35
    for idx, (num, name, body, color) in enumerate(steps):
        w = 1.72
        rect(s, x, 1.6, w, 1.85, fill=PANEL, line=PURPLE if idx == 4 else BORDER)
        pill(s, num, x + 0.14, 1.76, 0.42, color, size=11)
        add_text(s, name, x + 0.06, 2.26, w - 0.12, 0.5, size=11.3, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, body, x + 0.06, 2.82, w - 0.12, 0.55, size=9.8, color=TEXT, align=PP_ALIGN.CENTER)
        if idx < len(steps) - 1:
            arrow(s, x + w + 0.02, 2.45, 0.24, 0.3)
        x += w + 0.24
    rect(s, 0.45, 3.85, 6.1, 2.15, fill=PANEL, line=PURPLE)
    section_label(s, "왜 4단계(baseline)가 무료인가", 0.7, 4.1, 5.6, PURPLE)
    bullet_list(
        s,
        [
            "baseline은 AI 없이 미리 검증된 익스플로잇 명령으로 '타겟이 정말 취약한지'만 확인",
            "실패하면(=환경이 고장) 에이전트를 아예 실행하지 않고 종료 → 토큰 0",
            "'타겟 고장'과 '에이전트 실패'를 구분할 수 있게 해준다",
        ],
        0.7, 4.55, 5.6, 1.4, size=12, bullet_color=PURPLE,
    )
    rect(s, 6.75, 3.85, 6.13, 2.15, fill=PANEL, line=GREEN)
    section_label(s, "왜 정리(teardown)가 항상 보장되는가", 7.0, 4.1, 5.6, GREEN)
    bullet_list(
        s,
        [
            "성공·실패·강제 종료(Ctrl-C) 어떤 경우에도 finally 블록에서 실행",
            "compose down --volumes로 컨테이너·볼륨·네트워크까지 제거",
            "매 실행이 완전히 새 상태에서 시작 → 이전 실행이 다음에 영향 없음",
        ],
        7.0, 4.55, 5.6, 1.4, size=12, bullet_color=GREEN,
    )

    # 7. claude-loop engine -----------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "claude-loop 엔진 — AI가 지켜야 하는 규칙", "프롬프트로 '부탁'하는 게 아니라 시스템 프롬프트 + CLI로 강제한다", 7, TOTAL)
    card(s, "시스템 프롬프트 소유", "취약점은 '이미 확인됨'이라 알려줘\n정찰·스캔·열거 자체를 금지\n(nmap·gobuster·sqlmap 등 사용 안 함)", 0.45, 1.5, 6.0, 1.85, GREEN)
    card(s, "예산 상한 --max-budget-usd", "실행당 달러 하드 캡 (기본 $0.5)\nClaude CLI 자체가 강제 → 실제로\n예산을 넘기면 중단된다", 6.88, 1.5, 6.0, 1.85, YELLOW)
    card(s, "도구 화이트리스트", "기본은 Bash 도구 하나만 허용\n(--allow-web일 때만 WebFetch/\nWebSearch 추가) · 무인 실행", 0.45, 3.5, 6.0, 1.85, CYAN)
    card(s, "완료 프로토콜 + stop-on-proof", "성공하면 T9_RESULT 블록을 출력하고\n즉시 정지 · 그 후에도 도구를 계속\n쓰면(폭주) 강제 종료해 예산 절약", 6.88, 3.5, 6.0, 1.85, PURPLE)
    rect(s, 0.45, 5.5, 12.43, 1.0, fill=PANEL_2, line=CYAN)
    add_text(s, "참고", 0.7, 5.66, 1.0, 0.3, size=12, color=CYAN, bold=True)
    add_text(
        s,
        "ANTHROPIC_API_KEY는 사용하지 않는다 — 에이전트는 운영자의 Claude 구독(subscription)으로 실행되며, 대상은 지정된 host:port 하나뿐이다.",
        1.75, 5.66, 10.9, 0.7, size=12, color=WHITE,
    )

    # 8. Scenario file anatomy ---------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "시나리오 파일 구조 — 사람이 작성하는 설정", "scenarios/*.json 하나가 '무엇을 어떻게 공격할지'를 전부 정의한다", 8, TOTAL)
    code_block(
        s,
        '{\n'
        '  "t9_code": "T9-25-02-S-N-CD", "cve": "CVE-2017-5638",\n'
        '  "environment": { "path": "struts2/s2-045", "target_port": 8080 },\n'
        '  "agent": {\n'
        '    "budget_usd": 0.5, "proof": { "type": "reflected_http" },\n'
        '    "techniques": [ { "id": "classic-ognl-echo", "hint": "..." }, ... ],\n'
        '    "baseline": { "command": [...] }, "benign_profile": "..."\n'
        "  }\n}",
        0.45, 1.45, 6.0, 3.35, size=10.8,
    )
    section_label(s, "핵심 필드", 6.9, 1.4, 6.0, YELLOW)
    bullet_list(
        s,
        [
            "environment — 어느 Vulhub 경로/포트를 켤지, 준비 확인 방법(readiness)",
            "agent.techniques — 익스플로잇 기법 목록 (기법 뱅크, 다양성의 핵심)",
            "agent.proof — 성공을 어떻게 확인할지 (proof oracle 타입)",
            "agent.baseline — 무LLM 사전 점검용 정본 PoC 명령 (선택)",
            "agent.benign_profile — 정상 트래픽 샘플을 만들 때 줄 지침 (선택)",
            "agent.references — AI에게 줄 참고자료 (agentic RAG, 뒤에서 다룸)",
        ],
        6.9, 1.85, 6.0, 3.0, size=12.3, bullet_color=YELLOW,
    )
    add_text(
        s,
        "→ 사람은 '무엇을(CVE)·어떻게 판정할지·예산'만 정하고, 실제 payload 문자열은 에이전트가 실행 시점에 만든다",
        0.6, 5.0, 12.2, 0.4, size=12.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
    )

    # 9. Diversity engine -----------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "다양성 엔진 — 매번 다른 트래픽이 나오는 이유", "'다르게 해줘' 프롬프트가 아니라, 구조적으로 다양성을 강제한다", 9, TOTAL)
    card(s, "기법 뱅크 (Technique Bank)",
         "시나리오마다 서로 다른 익스플로잇\n기법 목록을 미리 정의\n실행마다 하나를 배정 → 그 자체가 라벨",
         0.45, 1.5, 6.0, 1.9, CYAN)
    card(s, "반복 방지 메모리",
         "직전 실행들이 어떤 기법을 썼는지\nmanifest에서 읽어 회피\n→ 같은 기법으로 수렴하지 않음",
         6.88, 1.5, 6.0, 1.9, GREEN)
    card(s, "모델 로테이션",
         "--repeat로 여러 번 돌리면 실행마다\n다른 모델을 사용 (models 목록)\n재시도 시 모델을 상향(escalation)",
         0.45, 3.6, 6.0, 1.9, PURPLE)
    card(s, "기법 재시도",
         "배정된 기법이 판정 실패하면\n(max_attempts까지) 다음 기법으로\n자동 전환 — 성공할 때까지 순환",
         6.88, 3.6, 6.0, 1.9, YELLOW)
    rect(s, 0.45, 5.72, 12.43, 0.75, fill=PANEL_2, line=CYAN)
    add_text(s, "예시", 0.7, 5.9, 1.0, 0.3, size=12, color=CYAN, bold=True)
    add_text(
        s,
        "Log4Shell 시나리오: JNDI in foo-param · JNDI in User-Agent 헤더 · alt-header/rmi — 같은 CVE, 3개의 서로 다른 주입점",
        1.55, 5.9, 11.2, 0.45, size=11.3, color=WHITE, bold=True,
    )

    # 10. Proof oracles -----------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "성공 판정 — 4가지 Proof Oracle", "성공은 에이전트의 말(transcript)이 아니라 'ground truth'로만 판정한다", 10, TOTAL)
    oracles = [
        ("reflected_http", "HTTP 응답 본문에 proof token이 등장 (반사형 RCE) — tshark로 응답 트래픽만 검사", GREEN),
        ("container_marker", "타깃이 토큰 이름의 파일을 실제로 생성 (blind RCE) — docker exec test -f", CYAN),
        ("container_log", "타깃 내부 로그 파일에 토큰이 기록됨 (예: Log4Shell) — docker exec grep", YELLOW),
        ("oob_callback", "캡처된 트래픽 어디에든 토큰이 등장 (아웃오브밴드 콜백) — 가장 약한 oracle", PURPLE),
    ]
    oy = 1.5
    for name, desc, color in oracles:
        rect(s, 0.45, oy, 12.43, 0.72, fill=PANEL, line=BORDER)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(oy), Inches(0.06), Inches(0.72))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_text(s, name, 0.66, oy + 0.2, 3.1, 0.32, size=13, color=color, bold=True, font=MONO)
        add_text(s, desc, 3.95, oy + 0.2, 8.75, 0.32, size=11.5, color=TEXT)
        oy += 0.82
    rect(s, 0.45, oy + 0.05, 12.43, 1.6, fill=PANEL_2, line=RED)
    section_label(s, "왜 에이전트의 말을 안 믿는가", 0.7, oy + 0.25, 6.0, RED)
    add_text(
        s,
        "curl -v 같은 명령은 요청 내용을 화면에 그대로 찍어준다. 그래서 '요청에 토큰을 넣었다'는 사실만으로도 에이전트가\n"
        "'성공했다'고 착각·과장 보고할 수 있다. proof oracle은 요청이 아니라 응답/서버 상태만 검사하므로, 실제로 침투가\n"
        "일어났는지와 무관한 '자기 신고'를 걸러낸다 — 에이전트가 성공을 주장해도 oracle이 아니라면 라벨은 '미확인'.",
        0.7, oy + 0.62, 11.9, 0.95, size=12, color=WHITE,
    )

    # 11. Baseline + retry (Tier 3) -------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "실패에 돈 쓰지 않기 — baseline 사전점검 + 기법 재시도", "AI를 실행하기 전에, AI 없이 먼저 확인한다", 11, TOTAL)
    rect(s, 0.48, 1.5, 6.0, 4.5, fill=PANEL, line=GREEN)
    section_label(s, "baseline 사전점검", 0.75, 1.76, 5.5, GREEN)
    bullet_list(
        s,
        [
            "시나리오에 정본(canonical) PoC 명령을 미리 등록해두면,",
            "에이전트를 실행하기 전에 그 명령을 그대로 실행해 본다 (LLM 없음, 완전 무료)",
            "실패하면 → '타겟 환경이 고장났다'로 판단하고 에이전트 실행을 건너뜀",
            "→ 애초에 뚫릴 수 없는 타겟에 토큰을 낭비하지 않는다",
        ],
        0.75, 2.24, 5.5, 3.4, size=13, bullet_color=GREEN,
    )
    rect(s, 6.9, 1.5, 6.0, 4.5, fill=PANEL, line=CYAN)
    section_label(s, "기법 재시도 (technique cycling)", 7.17, 1.76, 5.5, CYAN)
    bullet_list(
        s,
        [
            "배정된 기법으로 시도했는데 proof oracle이 실패를 판정하면,",
            "다음 기법으로 자동 전환해 재시도 (agent.max_attempts까지)",
            "재시도할 때마다 모델도 함께 상향할 수 있음 (retry_escalates_model)",
            "→ 한 기법이 막혀도 다른 방법으로 계속 시도, 각 시도는 개별 캡처로 저장",
        ],
        7.17, 2.24, 5.5, 3.4, size=13, bullet_color=CYAN,
    )
    add_text(
        s,
        "결과 = 뚫릴 수 없는 타겟에는 토큰을 쓰지 않고, 뚫릴 수 있는 타겟에는 성공할 때까지 다른 방법을 시도한다",
        0.6, 6.15, 12.2, 0.4, size=12.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
    )

    # 12. Exploit script + benign (Tier 4) ------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "산출물 — 재현 스크립트 + 정상 트래픽 샘플", "데이터셋에는 '공격'뿐 아니라 '정상'도 함께 담긴다", 12, TOTAL)
    rect(s, 0.48, 1.5, 6.0, 4.35, fill=PANEL, line=YELLOW)
    section_label(s, "exploit 스크립트 자동 저장", 0.75, 1.76, 5.5, YELLOW)
    bullet_list(
        s,
        [
            "익스플로잇에 성공하면 에이전트가 exploit.sh(또는 .py)를\n독립 실행 가능한 재현 스크립트로 직접 작성",
            "완벽하거나 범용적일 필요는 없음 — 읽을 수 있고 다시\n돌릴 수 있는 수준이면 충분",
            "이 시도의 실제 방법이 그대로 산출물에 남는다",
        ],
        0.75, 2.24, 5.5, 3.4, size=12.5, bullet_color=YELLOW,
    )
    rect(s, 6.9, 1.5, 6.0, 4.35, fill=PANEL, line=PURPLE)
    section_label(s, "정상 트래픽(benign) 샘플", 7.17, 1.76, 5.5, PURPLE)
    bullet_list(
        s,
        [
            "--benign 플래그로 실행하면 에이전트는 공격 대신\n'정상 사용자처럼' 행동 (일반 페이지 조회 등)",
            "같은 타겟·같은 하네스에서 얻은 '정상' 샘플이라\n탐지 모델의 hard-negative로 바로 쓸 수 있음",
            "proof oracle은 실행되지 않음 — 공격이 아니므로 판정 불필요",
        ],
        7.17, 2.24, 5.5, 3.4, size=12.5, bullet_color=PURPLE,
    )
    add_text(
        s,
        "각 샘플에는 기법·MITRE ATT&CK ID·주입점·payload 계열·비용·모델까지 함께 기록되어 학습에 바로 쓸 수 있다",
        0.6, 6.05, 12.2, 0.4, size=12.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
    )

    # 13. Agentic RAG -----------------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "Agentic RAG — 에이전트에게 힌트를 주는 두 계층", "모델이 모르는 CVE라도, 근거 자료를 주면 스스로 익스플로잇을 develop할 수 있다", 13, TOTAL)
    rect(s, 0.48, 1.5, 6.0, 4.15, fill=PANEL, line=GREEN)
    section_label(s, "오프라인 계층 (기본, 인터넷 불필요)", 0.75, 1.76, 5.5, GREEN)
    bullet_list(
        s,
        [
            "agent.references.notes — 사람이 적은 분석 메모(자유 텍스트)",
            "agent.references.source_paths — 타깃 컨테이너 내부 파일 경로\n(docker exec cat으로 하네스가 직접 읽어옴)",
            "에이전트는 docker 접근 권한이 없음 — 하네스가 대신 읽어서\n텍스트로 전달, 에이전트는 순수 네트워크 공격자로 유지",
        ],
        0.75, 2.24, 5.5, 3.2, size=12.3, bullet_color=GREEN,
    )
    rect(s, 6.9, 1.5, 6.0, 4.15, fill=PANEL, line=CYAN)
    section_label(s, "웹 계층 (--allow-web, 선택)", 7.17, 1.76, 5.5, CYAN)
    bullet_list(
        s,
        [
            "WebFetch/WebSearch 도구를 추가로 허용",
            "공개된 CVE의 공식 advisory·패치 diff를 읽고\n취약점의 원인·수정 지점을 이해하는 데 사용",
            "타깃을 정찰/스캔하는 데는 절대 쓸 수 없음 —\nCVE '지식' 조사 전용으로 범위가 제한됨",
        ],
        7.17, 2.24, 5.5, 3.2, size=12.3, bullet_color=CYAN,
    )
    rect(s, 0.48, 5.8, 12.4, 0.95, fill=PANEL_2, line=RED)
    add_text(s, "중요", 0.72, 5.96, 1.2, 0.3, size=12, color=RED, bold=True)
    add_text(
        s,
        "--allow-web은 반드시 '이미 공개된 CVE'에만 사용한다. 미공개·최신 취약점 조사에 쓰면 AI가 스스로 새 익스플로잇을\n찾아내도록 시키는 것과 같아 정책상 허용 범위를 벗어난다.",
        1.9, 5.96, 10.7, 0.7, size=11.5, color=WHITE,
    )

    # 14. Safety boundary --------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "안전 경계 — 이 시스템이 절대 하지 않는 것", "자유는 주되, 다섯 가지 원칙으로 범위를 강제한다", 14, TOTAL)
    rows = [
        ("공개된 CVE만", "모든 시나리오는 이미 알려지고 문서화된 CVE를 대상으로 한다 — 새 취약점을 찾는 도구가 아니다.", GREEN),
        ("격리된 일회용 타겟", "매 실행마다 새 Docker 컨테이너를 켜고, 끝나면 완전히 삭제 — 실제 서비스는 절대 건드리지 않는다.", CYAN),
        ("정찰·스캔 금지", "시스템 프롬프트가 nmap·gobuster·sqlmap 등 스캐너 사용을 명시적으로 금지한다.", YELLOW),
        ("예산·시간 상한", "--max-budget-usd(달러)와 timeout_seconds(초)가 모든 실행을 하드하게 제한한다.", PURPLE),
        ("자기 신고 불신", "성공 라벨은 항상 독립적인 proof oracle이 판정 — 에이전트의 주장은 참고용일 뿐이다.", RED),
    ]
    y = 1.5
    for name, desc, color in rows:
        rect(s, 0.45, y, 12.43, 0.88, fill=PANEL, line=BORDER)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(y), Inches(0.06), Inches(0.88))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_text(s, name, 0.66, y + 0.14, 2.6, 0.6, size=13, color=color, bold=True)
        add_text(s, desc, 3.35, y + 0.14, 9.35, 0.6, size=11.7, color=TEXT)
        y += 0.98

    # 15. Output directory -----------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "산출물 구조 — 폴더 안에 무엇이 저장되는가", "<T9 코드>/runs/<실행 ID>/ 아래에 모든 증거가 함께 남는다", 15, TOTAL)
    code_block(
        s,
        "<T9 코드>/runs/<실행 ID>/\n"
        "├── manifest.json         # 실행 메타데이터 + 최종 라벨(sample)\n"
        "├── capture.pcap          # 확정된 시도의 패킷 캡처\n"
        "├── agent_transcript.txt  # 에이전트의 전체 stream-json 로그\n"
        "├── exploit.sh            # 에이전트가 작성한 재현 스크립트\n"
        "├── baseline.txt          # 사전점검 명령의 출력 (있는 경우)\n"
        "├── attempts/a00/, a01/…  # 시도별 개별 캡처·트랜스크립트\n"
        "└── SHA256SUMS            # 모든 산출물의 체크섬",
        0.45, 1.45, 7.2, 3.3, size=11.5,
    )
    section_label(s, "manifest.json의 핵심", 7.85, 1.4, 5.0, YELLOW)
    bullet_list(
        s,
        [
            "engine · budget_usd · planned/used 기법 목록",
            "attempts — 시도별 기록(기법·모델·proof 결과)",
            "exploit_confirmed — ground truth 최종 판정",
            "sample — 학습에 바로 쓰는 라벨 블록\n(CVE·기법·MITRE ID·주입점·비용·모델)",
        ],
        7.85, 1.85, 4.9, 2.9, size=11.8, bullet_color=YELLOW,
    )
    rect(s, 0.45, 4.95, 12.43, 1.55, fill=PANEL_2, line=GREEN)
    add_text(s, "핵심", 0.7, 5.12, 1.0, 0.3, size=12, color=GREEN, bold=True)
    add_text(
        s,
        "manifest.json의 sample 블록이 실제 '제품' — 탐지 모델 학습에 바로 투입 가능한 라벨.\n"
        "SHA256SUMS 덕분에 나중에라도 산출물이 변조되지 않았음을 검증할 수 있다.",
        1.75, 5.12, 10.9, 1.2, size=12.3, color=WHITE,
    )

    # 16. Quick start + FAQ -------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "빠르게 시작하기", "체크리스트와 자주 하는 질문", 16, TOTAL)
    rect(s, 0.48, 1.5, 6.0, 4.6, fill=PANEL, line=CYAN)
    section_label(s, "준비물 체크리스트", 0.75, 1.76, 5.5, CYAN)
    bullet_list(
        s,
        [
            "Python 3.12 + uv (uv sync로 의존성 설치)",
            "Docker Compose 사용 권한",
            "~/vulhub 또는 $VULHUB_ROOT에 Vulhub 체크아웃",
            "Claude CLI 로그인(구독 계정) — API 키는 필요 없음",
            "선택: tshark/capinfos (캡처 확인용)",
        ],
        0.75, 2.24, 5.5, 3.6, size=13, bullet_color=CYAN,
    )
    rect(s, 6.9, 1.5, 6.0, 4.6, fill=PANEL, line=YELLOW)
    section_label(s, "자주 하는 질문", 7.17, 1.76, 5.5, YELLOW)
    bullet_list(
        s,
        [
            "Q. 비용은 얼마나 드나요?\nA. 실행당 --budget으로 정한 달러 상한 이내\n(기본 $0.5) — 이 예산을 넘기지 않는다.",
            "Q. 익스플로잇이 실패하면?\nA. proof oracle이 '미확인'으로 라벨링 —\n실패도 데이터로 남는다(정찰로 새지 않음).",
            "Q. 컨테이너가 남지 않나요?\nA. 성공·실패·중단 모두 finally에서\nteardown이 보장된다.",
        ],
        7.17, 2.24, 5.5, 3.6, size=12, bullet_color=YELLOW,
    )

    # 17. Summary + glossary -------------------------------------------------
    s = prs.slides.add_slide(blank)
    _bg(s)
    title(s, "요약 + 용어집", "핵심 개념을 한 문장씩", 17, TOTAL)
    terms = [
        ("claude-loop", "우리가 시스템 프롬프트·도구·예산을 완전히 통제하는 단일 claude -p 세션 엔진"),
        ("기법 뱅크 (technique bank)", "시나리오에 등록된 서로 다른 익스플로잇 방법 목록 — 다양성과 라벨의 원천"),
        ("proof token / oracle", "실행마다 발급되는 고유 토큰과, 그 토큰이 실제로 나타났는지 독립 검사하는 판정 로직"),
        ("ground truth", "에이전트의 주장이 아니라 실제 패킷·서버 상태로 판정하는 것 — 라벨 신뢰성의 핵심"),
        ("baseline PoC", "AI 없이 미리 검증하는 정본 익스플로잇 — 실패 시 토큰을 쓰지 않고 조기 종료"),
        ("benign hard-negative", "같은 타겟·하네스로 얻는 '정상' 트래픽 샘플 — 탐지 모델의 오탐 방지 학습에 사용"),
        ("agentic RAG", "에이전트에게 CVE 관련 참고자료(메모·소스·웹)를 주입해 모르는 취약점도 이해하게 돕는 기법"),
        ("teardown", "실행 성공/실패와 무관하게 컨테이너·네트워크를 항상 완전히 정리하는 절차"),
    ]
    y = 1.5
    for term, desc in terms:
        rect(s, 0.45, y, 12.43, 0.58, fill=PANEL, line=BORDER)
        add_text(s, term, 0.65, y + 0.1, 3.3, 0.38, size=12.3, color=CYAN, bold=True, font=MONO)
        add_text(s, desc, 4.05, y + 0.1, 8.6, 0.4, size=11.3, color=TEXT)
        y += 0.635
    add_text(
        s,
        "핵심 한 문장: 사람이 표적과 규칙을 고정하고, AI는 그 안에서 방법을 다양화하며, 결과는 항상 실제 증거로 검증된다.",
        0.6, 6.55, 12.2, 0.4, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
    )

    return prs


if __name__ == "__main__":
    presentation = build()
    presentation.save(OUT)
    print(OUT)
